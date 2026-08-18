#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Build the executive deck for "The Machine-Mediated Commerce Front".

Source : 2026.08.18-mcp-research-report.md
Output : 2026.08.18-Machine-Mediated-Commerce-Front-Executive-Deck.pptx

Design standard: Accenture brand identity system.
  - Accenture Purple #A100FF as the single accent, on black / white
  - Graphik is the brand typeface; Arial is Accenture's sanctioned Office fallback
  - Action titles (full-sentence "so what"), one idea per slide
  - Pyramid Principle: governing thought -> MECE sections -> evidence
"""

import os
from pptx import Presentation
from pptx.util import Inches as In, Pt
from pptx.dml.color import RGBColor as C
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------------------------------------------------------------- brand ----
PURPLE      = C(0xA1, 0x00, 0xFF)   # Accenture Purple (core)
PURPLE_DARK = C(0x75, 0x00, 0xC0)
PURPLE_DEEP = C(0x46, 0x00, 0x73)
PURPLE_LT   = C(0xBE, 0x82, 0xFF)
PURPLE_PALE = C(0xF0, 0xE4, 0xFF)
PURPLE_MIST = C(0xF9, 0xF4, 0xFF)

BLACK       = C(0x00, 0x00, 0x00)
INK         = C(0x1A, 0x1A, 0x24)
GREY        = C(0x6E, 0x6E, 0x78)
GREY_LT     = C(0xD9, 0xD9, 0xE0)
GREY_MIST   = C(0xF4, 0xF4, 0xF6)
WHITE       = C(0xFF, 0xFF, 0xFF)

FONT = "Arial"          # Graphik fallback per Accenture Office standards

# ------------------------------------------------------------ geometry ----
SW, SH   = 13.333, 7.5
M        = 0.55                 # side margin
CW       = SW - 2 * M           # content width
BODY_T   = 1.72                 # body top
BODY_B   = 6.02                 # body bottom when a takeaway band is present
BODY_B2  = 6.72                 # body bottom with no takeaway band

prs = Presentation()
prs.slide_width  = In(SW)
prs.slide_height = In(SH)
BLANK = prs.slide_layouts[6]

_n = [0]


# ------------------------------------------------------------- helpers ----
def slide():
    return prs.slides.add_slide(BLANK)


def rect(sl, x, y, w, h, fill=None, line=None, lw=0.75,
         shape=MSO_SHAPE.RECTANGLE, adj=None):
    s = sl.shapes.add_shape(shape, In(x), In(y), In(w), In(h))
    s.shadow.inherit = False
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(lw)
    if adj is not None:
        try:
            s.adjustments[0] = adj
        except Exception:
            pass
    tf = s.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = In(0.12)
    tf.margin_top = tf.margin_bottom = In(0.08)
    return s


def box(sl, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = sl.shapes.add_textbox(In(x), In(y), In(w), In(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    return tf


def para(tf, text, size=11, bold=False, color=INK, align=PP_ALIGN.LEFT,
         before=0, after=4, line=1.15, italic=False, first=False, caps=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(before)
    p.space_after = Pt(after)
    p.line_spacing = line
    r = p.add_run()
    r.text = text.upper() if caps else text
    f = r.font
    f.size, f.bold, f.italic, f.name = Pt(size), bold, italic, FONT
    f.color.rgb = color
    return p


def rich(tf, parts, size=11, align=PP_ALIGN.LEFT, before=0, after=4,
         line=1.15, first=False):
    """parts = [(text, bold, color) | (text, bold, color, size)]"""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(before)
    p.space_after = Pt(after)
    p.line_spacing = line
    for part in parts:
        txt, bold, col = part[0], part[1], part[2]
        sz = part[3] if len(part) > 3 else size
        r = p.add_run()
        r.text = txt
        r.font.size, r.font.bold, r.font.name = Pt(sz), bold, FONT
        r.font.color.rgb = col
    return p


def bullet(tf, text, size=10.5, color=INK, mark=PURPLE, after=5,
           line=1.18, first=False, bold=False):
    return rich(tf, [("▪  ", True, mark, size), (text, bold, color, size)],
                size=size, after=after, line=line, first=first)


def title_block(sl, kicker, title, dark=False):
    """Eyebrow + action title + Accenture rule."""
    tcol = WHITE if dark else INK
    tf = box(sl, M, 0.34, CW, 0.24)
    para(tf, kicker, size=8.5, bold=True, color=PURPLE, after=0,
         first=True, caps=True)
    tf = box(sl, M, 0.60, CW - 0.2, 0.82)
    para(tf, title, size=20.5, bold=True, color=tcol, after=0,
         line=1.06, first=True)
    rect(sl, M, 1.44, CW, 0.012, fill=GREY_LT if not dark else C(0x33, 0x33, 0x3D))
    rect(sl, M, 1.44, 1.15, 0.035, fill=PURPLE)


def takeaway(sl, text, y=6.06):
    rect(sl, M, y, CW, 0.60, fill=PURPLE_PALE)
    rect(sl, M, y, 0.045, 0.60, fill=PURPLE)
    tf = box(sl, M + 0.22, y + 0.09, CW - 0.42, 0.44, anchor=MSO_ANCHOR.MIDDLE)
    rich(tf, [("SO WHAT   ", True, PURPLE_DARK, 8.5),
              (text, False, INK, 10.5)],
         after=0, line=1.14, first=True)


def footer(sl, numbered=True):
    _n[0] += 1
    tf = box(sl, M, 6.94, 8.5, 0.2)
    para(tf, "Copyright © 2026 Accenture. All rights reserved."
             "   |   Confidential",
         size=7.5, color=GREY, after=0, first=True)
    if numbered:
        tf = box(sl, SW - M - 1.2, 6.94, 1.2, 0.2)
        para(tf, str(_n[0]), size=7.5, bold=True, color=PURPLE,
             align=PP_ALIGN.RIGHT, after=0, first=True)


def notes(sl, text):
    sl.notes_slide.notes_text_frame.text = text


def divider(num, title, sub, items):
    sl = slide()
    rect(sl, 0, 0, SW, SH, fill=BLACK)
    rect(sl, 0, 0, 0.09, SH, fill=PURPLE)

    tf = box(sl, M + 0.35, 1.62, 5.0, 1.7)
    para(tf, num, size=96, bold=True, color=PURPLE, after=0, line=0.92,
         first=True)

    tf = box(sl, M + 0.35, 3.30, 6.3, 1.5)
    para(tf, "SECTION", size=8.5, bold=True, color=GREY, after=6, first=True)
    para(tf, title, size=27, bold=True, color=WHITE, after=8, line=1.05)
    para(tf, sub, size=11, color=C(0x9A, 0x9A, 0xA6), after=0, line=1.25)

    rect(sl, 7.85, 1.95, 0.012, 3.6, fill=C(0x33, 0x33, 0x3D))
    tf = box(sl, 8.25, 1.95, 4.4, 3.6)
    para(tf, "IN THIS SECTION", size=8.5, bold=True, color=PURPLE,
         after=12, first=True)
    for it in items:
        bullet(tf, it, size=11, color=C(0xD5, 0xD5, 0xDC), mark=PURPLE,
               after=9, line=1.22)

    tf = box(sl, M + 0.35, 6.55, 2.0, 0.5)
    para(tf, ">", size=30, bold=True, color=PURPLE, after=0, first=True)
    _n[0] += 1
    tf = box(sl, SW - M - 1.2, 6.94, 1.2, 0.2)
    para(tf, str(_n[0]), size=7.5, bold=True, color=PURPLE,
         align=PP_ALIGN.RIGHT, after=0, first=True)
    return sl


def table(sl, x, y, w, headers, rows, colw, row_h=0.33, head_h=0.34,
          fsize=8.8, hsize=8.5, tint=None, chip_col=None, chip_map=None):
    """tint: set of row indices (0-based, body) to highlight."""
    tint = tint or set()
    nrows, ncols = len(rows) + 1, len(headers)
    gf = sl.shapes.add_table(nrows, ncols, In(x), In(y), In(w),
                             In(head_h + row_h * len(rows)))
    tbl = gf.table
    tbl.first_row = False
    tbl.horz_banding = False
    for i, cwid in enumerate(colw):
        tbl.columns[i].width = In(cwid)
    tbl.rows[0].height = In(head_h)
    for i in range(1, nrows):
        tbl.rows[i].height = In(row_h)

    def fill_cell(cell, txt, size, bold, color, bg, align=PP_ALIGN.LEFT):
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        cell.margin_left = cell.margin_right = In(0.07)
        cell.margin_top = cell.margin_bottom = In(0.03)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        p.line_spacing = 1.0
        p.space_after = Pt(0)
        r = p.add_run()
        r.text = txt
        r.font.size, r.font.bold, r.font.name = Pt(size), bold, FONT
        r.font.color.rgb = color

    for j, h in enumerate(headers):
        fill_cell(tbl.cell(0, j), h.upper(), hsize, True, WHITE, PURPLE_DEEP)

    for i, row in enumerate(rows):
        bg = PURPLE_PALE if i in tint else (WHITE if i % 2 == 0 else GREY_MIST)
        for j, val in enumerate(row):
            fg, bold, cbg = INK, (j == 0), bg
            if chip_col is not None and j == chip_col and chip_map:
                cfill, ctext = chip_map.get(val, (bg, INK))
                cbg, fg, bold = cfill, ctext, True
            fill_cell(tbl.cell(i + 1, j), val, fsize, bold, fg, cbg)
    return tbl


CHIPS = {
    "No-regret foundation": (PURPLE_DEEP, WHITE),
    "Near-term priority":   (PURPLE,      WHITE),
    "Strategic option":     (PURPLE_LT,   PURPLE_DEEP),
    "Frontier bet":         (PURPLE_PALE, PURPLE_DARK),
    "Monitor":              (PURPLE_PALE, PURPLE_DARK),
    "Avoid":                (GREY_LT,     GREY),
    "Avoid / monitor":      (GREY_LT,     GREY),
}


def card(sl, x, y, w, h, eyebrow, heading, lines, accent=PURPLE,
         fill=WHITE, border=GREY_LT, hsize=12, bsize=9.5, eyebrow_size=8):
    s = rect(sl, x, y, w, h, fill=fill, line=border, lw=0.75)
    s.text_frame.text = ""
    rect(sl, x, y, w, 0.035, fill=accent)
    tf = box(sl, x + 0.20, y + 0.24, w - 0.40, h - 0.40)
    first = True
    if eyebrow:
        para(tf, eyebrow, size=eyebrow_size, bold=True, color=accent,
             after=5, first=True, caps=True)
        first = False
    if heading:
        para(tf, heading, size=hsize, bold=True, color=INK, after=7,
             line=1.10, first=first)
        first = False
    for ln in lines:
        bullet(tf, ln, size=bsize, color=C(0x3A, 0x3A, 0x46), mark=accent,
               after=5, line=1.18, first=first)
        first = False
    return s


def spoke(sl, x1, y1, x2, y2, color=PURPLE_LT, lw=1.0):
    """Thin connector, drawn before nodes so the nodes sit on top of it."""
    c = sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                In(x1), In(y1), In(x2), In(y2))
    c.line.color.rgb = color
    c.line.width = Pt(lw)
    return c


def vertical(tf):
    """Rotate a text frame 270° for use as a vertical axis label."""
    tf._txBody.bodyPr.set("vert", "vert270")
    return tf


def statile(sl, x, y, w, h, figure, label, source, accent=PURPLE):
    rect(sl, x, y, w, h, fill=PURPLE_MIST, line=GREY_LT, lw=0.75)
    rect(sl, x, y, 0.04, h, fill=accent)
    tf = box(sl, x + 0.24, y + 0.26, w - 0.44, h - 0.45)
    para(tf, figure, size=34, bold=True, color=accent, after=4, line=1.0,
         first=True)
    para(tf, label, size=10, bold=True, color=INK, after=5, line=1.14)
    para(tf, source, size=7.5, color=GREY, after=0, line=1.14)


# ============================================================ SLIDE  1 ====
sl = slide()
rect(sl, 0, 0, SW, SH, fill=BLACK)
rect(sl, 0, 0, SW, 0.10, fill=PURPLE)
rect(sl, 8.55, 0.10, 4.78, SH - 0.10, fill=C(0x0D, 0x0D, 0x12))
tf = box(sl, 9.05, 5.72, 3.9, 0.9)
para(tf, ">", size=54, bold=True, color=PURPLE, align=PP_ALIGN.RIGHT,
     after=0, first=True)

tf = box(sl, M + 0.35, 1.55, 7.6, 0.3)
para(tf, "Strategic research programme   |   Research cutoff 18 August 2026",
     size=9.5, bold=True, color=PURPLE, after=0, first=True, caps=True)

tf = box(sl, M + 0.35, 2.05, 7.7, 2.2)
para(tf, "The Machine-Mediated\nCommerce Front", size=42, bold=True,
     color=WHITE, after=0, line=1.02, first=True)

rect(sl, M + 0.35, 4.32, 1.5, 0.045, fill=PURPLE)

tf = box(sl, M + 0.35, 4.62, 7.4, 1.3)
para(tf, "Where value actually accrues as agents — not people — become "
         "the buyers, and what the firm should advise, build and invest in "
         "across three horizons.",
     size=13, color=C(0xB4, 0xB4, 0xBE), after=0, line=1.32, first=True)

tf = box(sl, M + 0.35, 6.55, 7.4, 0.4)
para(tf, "Prepared for the executive team of a multi-domain consulting firm",
     size=9, color=GREY, after=0, first=True)

tf = box(sl, 9.05, 2.05, 3.9, 3.2)
para(tf, "AT A GLANCE", size=8.5, bold=True, color=PURPLE, after=14,
     first=True)
for k, v in [("17%", "of shoppers are comfortable buying through AI"),
             ("41%", "of software organisations have MCP in production"),
             ("88%", "reported an AI-agent security incident this year"),
             ("3–7 yrs", "for the economic restructuring to play out")]:
    rich(tf, [(k, True, WHITE, 20)], after=1, line=1.0)
    para(tf, v, size=8.5, color=GREY, after=13, line=1.2)

notes(sl, "One-line frame: the shift is real, but the popular narrative has "
          "the timing and the location of value wrong. This deck resolves five "
          "questions in sequence.")
_n[0] += 1

# ============================================================ SLIDE  2 ====
sl = slide()
title_block(sl, "How to read this deck",
            "This deck answers five questions in sequence — each section "
            "closes one before the next opens")
qs = [("01", "Is the shift real, and is it on time?",
       "The shift is real. The popular timing and location of value are wrong."),
      ("02", "Where does value actually accrue?",
       "To scarce control points — not to protocols, which are commoditising."),
      ("03", "When does it land, and in which sectors?",
       "Enterprise and B2B first; consumer autonomy last. Six sectors ranked."),
      ("04", "What do we do about it?",
       "Three separate theses: for clients, for the firm, for external investment."),
      ("05", "What would prove us wrong?",
       "Five falsifiers and seven leading indicators, declared up front.")]
y = BODY_T
for num, q, a in qs:
    rect(sl, M, y, CW, 0.74, fill=WHITE, line=GREY_LT, lw=0.75)
    rect(sl, M, y, 0.9, 0.74, fill=PURPLE_DEEP)
    tf = box(sl, M, y + 0.17, 0.9, 0.4)
    para(tf, num, size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
         after=0, first=True)
    tf = box(sl, M + 1.15, y + 0.13, 5.1, 0.48, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, q, size=11.5, bold=True, color=INK, after=0, line=1.1, first=True)
    tf = box(sl, M + 6.45, y + 0.13, CW - 6.65, 0.48, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, a, size=9.5, color=GREY, after=0, line=1.18, first=True)
    y += 0.82
takeaway(sl, "Each section is mutually exclusive and collectively exhaustive: "
             "no argument appears twice, and nothing material is left out.")
footer(sl)
notes(sl, "Set expectations: this is a decision deck, not a literature review. "
          "Every section ends with a decision or a declared uncertainty.")

# ============================================================ SLIDE  3 ====
sl = slide()
title_block(sl, "The executive answer",
            "The working thesis is confirmed — with one correction that "
            "changes where we invest first")

card(sl, M, BODY_T, 5.95, 3.05, "Substantially confirmed",
     "Machine-mediated commerce is real, accelerating and structurally "
     "significant",
     ["Commerce is moving from human interfaces to machine-readable "
      "capabilities — high confidence",
      "Value accrues to controllers of intent, data, capability, identity, "
      "orchestration, settlement and assurance",
      "Open, foundation-governed standards are winning the protocol layer",
      "The restructuring is a three-to-seven year phenomenon, not a "
      "this-year event"],
     accent=PURPLE, hsize=13, bsize=10)

card(sl, M + 6.28, BODY_T, 5.95, 3.05, "One material correction",
     "The Horizon-1 front line is enterprise and the rails — not "
     "autonomous consumer shopping",
     ["Value is in enterprise back-office and B2B workflows: procurement, "
      "service, finance, compliance",
      "And in the payment, identity and trust “rails” that make "
      "delegation safe",
      "Autonomous consumer checkout stalled in early 2026 on trust and "
      "conversion economics",
      "It is a Horizon-2/3 phenomenon — fund it as an option, not as a "
      "plan"],
     accent=PURPLE_DEEP, fill=PURPLE_MIST, hsize=13, bsize=10)

takeaway(sl, "Build around the durable control points — not around any "
             "single protocol, and not around the consumer-checkout narrative. "
             "MCP is now infrastructure, but it is one layer of a "
             "multi-protocol stack.")
footer(sl)
notes(sl, "The governing thought of the whole deck. If the audience remembers "
          "one slide, it is this one. The correction is the commercially "
          "consequential part: it redirects Horizon-1 spend.")

# ============================================================ SLIDE  4 ====
sl = slide()
title_block(sl, "The correction, in evidence",
            "Consumers accept agents as shopping assistants far faster than "
            "as autonomous buyers")

ev = [("March 2026", "OpenAI scaled back in-chat Instant Checkout",
       "Only ~12 of Shopify’s millions of merchants had gone live. The "
       "single most decisive datapoint of 2026."),
      ("Walmart", "In-chat checkout converted ~3× worse",
       "The same shoppers sent to walmart.com converted at roughly three times "
       "the rate of completing in chat."),
      ("February 2026", "eBay banned unauthorised buy-for-me bots",
       "Marketplaces are actively fencing off autonomous purchasing rather "
       "than courting it.")]
x = M
for date, head, body in ev:
    card(sl, x, BODY_T, 3.94, 2.10, date, head, [body],
         accent=PURPLE, hsize=11.5, bsize=9.5)
    x += 4.14

rect(sl, M, BODY_T + 2.32, CW, 1.32, fill=BLACK)
rect(sl, M, BODY_T + 2.32, 0.045, 1.32, fill=PURPLE)
tf = box(sl, M + 0.30, BODY_T + 2.50, 5.0, 1.0)
para(tf, "17%", size=40, bold=True, color=PURPLE, after=0, line=1.0,
     first=True)
tf = box(sl, M + 1.85, BODY_T + 2.56, 10.0, 1.0)
para(tf, "of shoppers feel comfortable completing a purchase through AI",
     size=13.5, bold=True, color=WHITE, after=4, line=1.1, first=True)
para(tf, "ChannelEngine Marketplace Shopping Behavior Report 2026. Trust — "
         "not model capability — is what is binding.",
     size=9, color=C(0x9A, 0x9A, 0xA6), after=0, line=1.18)

takeaway(sl, "“Discover in AI, buy on site” is the durable pattern "
             "today. Advise clients to fund machine-readable capability and "
             "governance — not chatbot checkout.")
footer(sl)
notes(sl, "Three independent signals from three different actor types "
          "(assistant provider, retailer, marketplace) all pointing the same "
          "way. That convergence is what makes the correction safe to act on.")

# ============================================================ SLIDE  5 ====
sl = slide()
title_block(sl, "The core architecture",
            "Machine-mediated commerce nests in three scopes — and "
            "Horizon-1 value sits in the middle ring, not the inner one")

# --- three concentric scopes -------------------------------------------
rect(sl, 0.55, 1.72, 7.30, 3.28, fill=PURPLE_MIST, line=PURPLE_LT, lw=1.0,
     shape=MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.10)
tf = box(sl, 0.82, 1.84, 6.76, 0.44)
rich(tf, [("1   ", True, PURPLE_LT, 12),
          ("MACHINE-MEDIATED ECONOMIC INTERACTION", True, PURPLE_DARK, 9.5)],
     after=2, line=1.0, first=True)
para(tf, "Autonomous exchanges between agents, machines, organisations and "
         "digital services", size=8.2, color=GREY, after=0, line=1.12)

rect(sl, 1.05, 2.30, 6.30, 2.22, fill=WHITE, line=PURPLE, lw=1.0,
     shape=MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.12)
tf = box(sl, 1.32, 2.42, 5.76, 0.44)
rich(tf, [("2   ", True, PURPLE_LT, 12),
          ("ENTERPRISE COMMERCE", True, PURPLE_DARK, 9.5)],
     after=2, line=1.0, first=True)
para(tf, "Procurement, sales, supply chain and finance, executed via agents",
     size=8.2, color=GREY, after=0, line=1.12)

rect(sl, 1.55, 2.88, 5.30, 1.16, fill=PURPLE_PALE, line=PURPLE_DEEP, lw=1.25,
     shape=MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.16)
tf = box(sl, 1.82, 3.02, 4.76, 0.90)
rich(tf, [("3   ", True, PURPLE, 12),
          ("CUSTOMER COMMERCE", True, PURPLE_DEEP, 9.5)],
     after=2, line=1.0, first=True)
para(tf, "Discovery, purchase, payment and fulfilment via personal agents",
     size=8.2, color=C(0x3A, 0x3A, 0x46), after=3, line=1.12)
para(tf, "The smallest ring — and the one that stalled in 2026",
     size=8.2, bold=True, color=PURPLE_DEEP, after=0, line=1.12)

# --- the trust substrate beneath all three -----------------------------
rect(sl, 4.03, 5.02, 0.34, 0.21, fill=PURPLE, shape=MSO_SHAPE.UP_ARROW)
rect(sl, 0.55, 5.25, 7.30, 0.60, fill=BLACK)
rect(sl, 0.55, 5.25, 0.045, 0.60, fill=PURPLE)
tf = box(sl, 0.85, 5.34, 6.85, 0.44, anchor=MSO_ANCHOR.MIDDLE)
rich(tf, [("TRUST INFRASTRUCTURE   ", True, PURPLE, 8.5),
          ("Identity · delegated authority · consent · policy · security · "
           "liability · reversibility", False, WHITE, 8.8)],
     after=0, line=1.14, first=True)
tf = box(sl, 4.45, 4.99, 3.3, 0.24)
para(tf, "constrains and enables every ring above", size=7.5, italic=True,
     color=GREY, after=0, first=True)

# --- what MCP does not solve -------------------------------------------
rect(sl, 8.20, 1.72, 4.58, 4.13, fill=WHITE, line=GREY_LT, lw=0.75)
rect(sl, 8.20, 1.72, 4.58, 0.035, fill=PURPLE_DEEP)
tf = box(sl, 8.46, 1.96, 4.06, 3.70)
para(tf, "MCP LOWERS THE COST OF ACCESS", size=8.5, bold=True,
     color=PURPLE_DARK, after=4, first=True)
para(tf, "But it does not solve:", size=11.5, bold=True, color=INK, after=11,
     line=1.10)
for it, where in [("Agent discovery", "layer 1"),
                  ("Agent-to-agent coordination", "A2A, layer 3"),
                  ("Commerce semantics", "ACP / UCP, layer 5"),
                  ("Payment and settlement", "AP2 / x402, layer 5"),
                  ("Liability and governance", "layer 6")]:
    rich(tf, [("▪  ", True, PURPLE, 9.5), (it, True, INK, 9.5)],
         after=1, line=1.14)
    para(tf, "        " + where, size=8, color=GREY, after=8, line=1.10)

takeaway(sl, "This is why the deck is organised around control points rather "
             "than around MCP. Each ring needs a different protocol, and all "
             "three rest on a trust layer that no protocol supplies.")
footer(sl)
notes(sl, "Use the rings to place the correction physically: Horizon-1 value "
          "is the middle ring. The inner ring is the smallest, hardest and "
          "the one that stalled. The trust bar underneath is section 01's "
          "argument in one shape.")

# ============================================================ SLIDE  6 ====
sl = slide()
title_block(sl, "The argument in three propositions",
            "Three propositions carry the whole argument — everything "
            "that follows is evidence for one of them")

props = [("01", "The shift is real, but the popular narrative is mis-timed",
          ["Enterprise and B2B agentic workflows plus the payment-identity-"
           "trust rails are the Horizon-1 value",
           "Autonomous consumer shopping stalled on trust and conversion "
           "economics in early 2026",
           "Advise clients to invest in machine-readable capabilities and "
           "governance, not chatbot checkout"]),
         ("02", "Value accrues to control points, not protocols",
          ["MCP, A2A, ACP, UCP, AP2 and x402 are all converging on open, "
           "foundation governance",
           "That convergence will substantially commoditise them",
           "Defensible value sits in identity and delegated authority, "
           "differentiated data, orchestration, settlement and assurance"]),
         ("03", "For the firm, this is a no-regret foundation play — now",
          ["Build a proprietary maturity diagnostic, a trust/control library "
           "and protocol accelerators",
           "Stand up a managed agent-operations and assurance service",
           "Partner — do not build — on observability and payments; "
           "treat autonomous operators as a frontier bet"])]
x = M
for num, head, lines in props:
    rect(sl, x, BODY_T, 3.94, 4.05, fill=WHITE, line=GREY_LT, lw=0.75)
    rect(sl, x, BODY_T, 3.94, 0.035, fill=PURPLE)
    tf = box(sl, x + 0.22, BODY_T + 0.24, 3.5, 0.55)
    para(tf, num, size=27, bold=True, color=PURPLE_LT, after=0, first=True)
    tf = box(sl, x + 0.22, BODY_T + 0.86, 3.5, 3.0)
    para(tf, head, size=12.5, bold=True, color=INK, after=11, line=1.10,
         first=True)
    for ln in lines:
        bullet(tf, ln, size=9.5, color=C(0x3A, 0x3A, 0x46), after=7, line=1.20)
    x += 4.14

takeaway(sl, "Sections 01–03 establish these propositions; sections "
             "04–05 convert them into decisions and declare what would "
             "overturn them.")
footer(sl)
notes(sl, "This is the pyramid. Each proposition maps to one body section. "
          "Do not add a fourth — the discipline of three is what keeps "
          "the argument decision-ready.")

# ========================================================== DIVIDER 01 ====
divider("01", "The shift is real —\nbut mis-timed",
        "The infrastructure arrived in 2026. The trust did not.",
        ["The interface shift: from attracting humans to serving agents",
         "Four milestones that made the infrastructure real",
         "Why announcements are not adoption",
         "Trust as the binding constraint on autonomy",
         "Where Horizon-1 value therefore sits",
         "The ten gates every autonomous use case must pass"])

# ============================================================ SLIDE  7 ====
sl = slide()
title_block(sl, "The interface shift",
            "The commercial question changes from “how do we attract a "
            "human to our interface?” to “how does a trusted agent "
            "transact with our capabilities?”")

PH = 3.25
rect(sl, M, BODY_T, 5.45, PH, fill=GREY_MIST, line=GREY_LT, lw=0.75)
tf = box(sl, M + 0.28, BODY_T + 0.24, 4.9, 0.4)
para(tf, "TRADITIONAL COMMERCE", size=10, bold=True, color=GREY, after=0,
     first=True)
y = BODY_T + 0.70
for item in ["Pages and apps", "Funnels", "Search rankings",
             "Advertising", "Call centres", "Manual workflows"]:
    rect(sl, M + 0.28, y, 4.9, 0.36, fill=WHITE, line=GREY_LT, lw=0.75)
    tf = box(sl, M + 0.48, y + 0.08, 4.5, 0.24)
    para(tf, item, size=10, color=GREY, after=0, first=True)
    y += 0.38

rect(sl, M + 5.72, BODY_T + 1.32, 0.90, 0.62, fill=PURPLE,
     shape=MSO_SHAPE.RIGHT_ARROW)

rect(sl, M + 6.88, BODY_T, 5.45, PH, fill=PURPLE_MIST, line=PURPLE_LT,
     lw=0.75)
tf = box(sl, M + 7.16, BODY_T + 0.24, 4.9, 0.4)
para(tf, "AGENT-MEDIATED COMMERCE", size=10, bold=True, color=PURPLE_DARK,
     after=0, first=True)
y = BODY_T + 0.70
for item in ["Structured intent", "Machine-readable capabilities",
             "Verifiable data", "Policy-compliant actions",
             "Measurable outcomes", "Agent discovery and selection"]:
    rect(sl, M + 7.16, y, 4.9, 0.36, fill=WHITE, line=PURPLE_LT, lw=0.75)
    tf = box(sl, M + 7.36, y + 0.08, 4.5, 0.24)
    para(tf, item, size=10, bold=True, color=INK, after=0, first=True)
    y += 0.38

# the framing question each model asks
for qx, qtext, col, bg in [
        (M, "How do we attract a human to our interface?", GREY, GREY_MIST),
        (M + 6.88, "How does a trusted agent discover and transact with our "
                   "capabilities?", PURPLE_DARK, PURPLE_PALE)]:
    rect(sl, qx, 5.08, 5.45, 0.80, fill=bg)
    rect(sl, qx, 5.08, 0.045, 0.80, fill=col)
    tf = box(sl, qx + 0.26, 5.18, 5.0, 0.60, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, "“" + qtext + "”", size=10.5, bold=True, color=col,
         after=0, line=1.20, first=True)

takeaway(sl, "An 8,000-vendor martech stack is built on the assumption that "
             "buyers arrive at seller environments. Agents break that "
             "assumption — and with it, the discovery and SEO value pool.")
footer(sl)
notes(sl, "Frame this as an assumption failure, not a channel change. Everything "
          "downstream — martech spend, SEO, retail media — inherits "
          "the broken assumption.")

# ============================================================ SLIDE  8 ====
sl = slide()
title_block(sl, "Infrastructure",
            "In 2026 the infrastructure layer stopped being experimental — "
            "four milestones made it real")

ms = [("28 July 2026", "MCP reached its stateless specification",
       ["OAuth 2.1 / OIDC aligned; Resource Indicators (RFC 8707) required",
        "Enterprise-Managed Authorization now stable — Anthropic, "
        "Microsoft, Okta",
        "~half a billion monthly Tier-1 SDK downloads; ~9,650 registry servers"]),
      ("9 April 2026", "A2A reached v1.0 — production-grade",
       ["Signed Agent Cards and multi-tenancy for cross-boundary coordination",
        "150+ organisations; Linux Foundation governed, absorbed IBM ACP",
        "Shipped in Vertex AI, Bedrock AgentCore, Azure AI Foundry"]),
      ("26 May 2026", "Payment authority consolidated at FIDO",
       ["Google’s AP2 and Mastercard’s Verifiable Intent both donated",
        "AP2 v0.2 adds “Human Not Present”; 60+ partners",
        "Payments TWG chaired by Mastercard and Visa"]),
      ("Through 2026", "Card networks shipped production agentic rails",
       ["Visa Intelligent Commerce Connect unifies four protocols",
        "Mastercard Agent Pay live in Hong Kong (Mar) and Thailand (Apr)",
        "x402 became a Linux Foundation project, operational 14 July 2026"])]
x, y = M, BODY_T
for i, (date, head, lines) in enumerate(ms):
    card(sl, x, y, 5.95, 1.72, date, head, lines, accent=PURPLE,
         hsize=11.5, bsize=8.8)
    x = M + 6.28 if i % 2 == 0 else M
    if i % 2 == 1:
        y += 1.90

takeaway(sl, "The plumbing is no longer the constraint. Any strategy that is "
             "still waiting for the standards to settle is now a year behind.")
footer(sl)
notes(sl, "Four different governance bodies, four different layers, all within "
          "twelve months. That is what ‘de facto standard’ looks like "
          "in practice.")

# ============================================================ SLIDE  9 ====
sl = slide()
title_block(sl, "Adoption reality",
            "Announcements are not adoption — production reality is "
            "roughly half of what the headlines imply")

stats = [("41%", "of surveyed software organisations have MCP in limited or "
                 "broad production", "Stacklok 2026 software report"),
         ("8.5%", "of MCP servers implement mandatory OAuth 2.1",
          "The security floor is not being met"),
         ("14.4%", "of agents went live with full security and IT approval",
          "Gravitee State of AI Agent Security 2026"),
         ("88%", "of organisations reported a confirmed or suspected AI-agent "
                 "security incident", "Gravitee 2026; healthcare 92.7%")]
x = M
for fig, lab, src in stats:
    statile(sl, x, BODY_T, 2.92, 2.05, fig, lab, src)
    x += 3.11

rect(sl, M, BODY_T + 2.28, CW, 1.36, fill=WHITE, line=PURPLE_LT, lw=1.0)
rect(sl, M, BODY_T + 2.28, 0.045, 1.36, fill=PURPLE)
tf = box(sl, M + 0.30, BODY_T + 2.46, CW - 0.6, 1.05)
para(tf, "THREE HEADLINE NUMBERS THAT DO NOT SURVIVE TRIANGULATION",
     size=9, bold=True, color=PURPLE_DARK, after=8, first=True)
bullet(tf, "“MCP is 78% in production” is debunked — use ~41%",
       size=9.5, after=4)
bullet(tf, "x402 cumulative payment counts do not reconcile: Coinbase figures "
           "range ~75M to ~169M, on only $24–50M cumulative value",
       size=9.5, after=4)
bullet(tf, "AgentCore “~200ms settlement” is the settlement step "
           "only — independent testing measured ~3–5s end to end",
       size=9.5, after=0)

takeaway(sl, "Read every vendor number as a pilot number until proven "
             "otherwise. Our credibility with clients depends on being the "
             "firm that corrects these, not the one that repeats them.")
footer(sl)
notes(sl, "This slide is a trust-builder with the executive audience. It says "
          "we did the triangulation. Use it to pre-empt ‘but I read "
          "that...’ challenges.")

# =========================================================== SLIDE 10 ====
sl = slide()
title_block(sl, "The binding constraint",
            "Trust, not technology, is what stops autonomy — and trust is "
            "an engineering programme, not a communications campaign")

card(sl, M, BODY_T, 3.94, 2.55, "The demand-side deficit",
     "Users will not delegate",
     ["Only 17% of shoppers are comfortable buying through AI",
      "eBay banned unauthorised buy-for-me bots; Walmart saw ~3× worse "
      "conversion",
      "Liability for agent errors is unresolved across the FTC, the FCA "
      "and PSD3"], accent=PURPLE_DEEP, hsize=11.5, bsize=9.2)

card(sl, M + 4.14, BODY_T, 3.94, 2.55, "The supply-side attack surface",
     "Agents introduce new classes of attack",
     ["Tool poisoning, descriptor poisoning, agent-in-the-middle, "
      "goal hijacking",
      "Documented in the OWASP MCP Top 10 and peer-reviewed research",
      "Real exploitation: CVE-2025-59528 in Flowise; a malicious MCP server "
      "impersonating Postmark"], accent=PURPLE_DEEP, hsize=11.5, bsize=9.2)

card(sl, M + 8.28, BODY_T, 3.95, 2.55, "The unlock",
     "Five controls turn autonomy on",
     ["Just-in-time ephemeral scoped tokens",
      "Human-in-the-loop for every state-mutating operation",
      "Signed manifests and agent cards; allowlists",
      "Continuous monitoring and audit evidence"],
     accent=PURPLE, fill=PURPLE_MIST, hsize=11.5, bsize=9.2)

rect(sl, M, BODY_T + 2.78, CW, 0.86, fill=BLACK)
tf = box(sl, M + 0.30, BODY_T + 2.96, CW - 0.6, 0.55, anchor=MSO_ANCHOR.MIDDLE)
rich(tf, [("The strategic implication   ", True, PURPLE, 10),
          ("Every consulting engagement in Horizon 1 is, underneath, a trust "
           "engineering engagement. That is a services business — not a "
           "software business — and it is ours to take.",
           False, WHITE, 11)], after=0, line=1.2, first=True)

takeaway(sl, "88% incident rate plus 14.4% approval rate is the commercial "
             "opening: assurance is the highest willingness-to-pay category "
             "in the stack.")
footer(sl)
notes(sl, "Pivot point of section 01. The constraint is also the opportunity: "
          "what blocks clients is exactly what they will pay us to fix.")

# =========================================================== SLIDE 11 ====
sl = slide()
title_block(sl, "Where Horizon-1 value sits",
            "Because trust gates autonomy, Horizon-1 value sits where a human "
            "still approves — enterprise and B2B workflows")

rows = [
    ["Procurement research and shortlisting", "Discovery → evaluation",
     "51% of B2B software buyers now begin in an AI chatbot; 69% chose a "
     "different vendor than planned", "Recommendatory", "Low"],
    ["Quote-to-cash and PO matching", "Authorisation → ordering",
     "Up to 80% cycle-time reduction in PO processing (PwC); 40–60% "
     "cycle cuts (Salesforce)", "Approval-based", "Low"],
    ["Credit risk memos", "Evaluation → decision",
     "20–60% productivity gain and 30% credit-turnaround improvement at "
     "a US bank (McKinsey)", "Approval-based", "Low"],
    ["Financial crime and compliance", "Service → renewal",
     "FIS + Anthropic Financial Crimes AI Agent; BMO and Amalgamated Bank as "
     "early deployers", "Governed, auditable", "Low"],
    ["Inventory monitoring and replenishment", "Ordering → fulfilment",
     "Agents monitor stock, compare terms and trigger purchases inside "
     "existing contracts", "Bounded autonomy", "Medium"],
    ["Customer service and disruption recovery", "Service",
     "Malaysia Airlines “Mavis” end-to-end service agent; Navan "
     "travel rebooking", "Bounded autonomy", "Medium"],
    ["Consumer discovery and comparison", "Discovery → evaluation",
     "High adoption — but the transaction step stalls on trust",
     "Recommendatory", "High at checkout"],
]
table(sl, M, BODY_T, CW,
      ["Use case", "Lifecycle stage", "Evidence of value", "Autonomy level",
       "Exposure"],
      rows, [3.05, 1.85, 4.55, 1.55, 1.23], row_h=0.44, fsize=8.6,
      tint={0, 1, 2, 3})

takeaway(sl, "The four highlighted use cases combine quantified production "
             "outcomes with low exposure. They are where the first engagement "
             "should start — in every sector.")
footer(sl)
notes(sl, "Note the pattern: every low-exposure case keeps a human on the "
          "authorisation step. Autonomy is earned stage by stage, not granted.")

# =========================================================== SLIDE 12 ====
sl = slide()
title_block(sl, "The autonomy gates",
            "Ten gates must all pass before a use case may execute "
            "autonomously — and they cannot be averaged away")

gates = ["Identified actor and owner", "Authenticated identity",
         "Delegated authority", "Data rights and permissions",
         "Secure capability access", "Transaction limits and controls",
         "Human intervention and reversibility", "Evaluation and monitoring",
         "Audit evidence and incident response",
         "Legal review and acceptable unit economics"]
x, y = M, BODY_T
for i, g in enumerate(gates):
    rect(sl, x, y, 2.38, 1.24, fill=WHITE, line=PURPLE_LT, lw=0.9)
    rect(sl, x, y, 2.38, 0.035, fill=PURPLE)
    tf = box(sl, x + 0.16, y + 0.22, 2.06, 0.22)
    para(tf, "GATE %02d" % (i + 1), size=7.5, bold=True, color=PURPLE_LT,
         after=0, first=True)
    tf = box(sl, x + 0.16, y + 0.54, 2.06, 0.58)
    para(tf, g, size=9.5, bold=True, color=INK, after=0, line=1.12, first=True)
    x += 2.47
    if i == 4:
        x, y = M, y + 1.40

rect(sl, M, BODY_T + 3.06, CW, 0.92, fill=PURPLE_DEEP)
tf = box(sl, M + 0.30, BODY_T + 3.24, CW - 0.6, 0.56, anchor=MSO_ANCHOR.MIDDLE)
rich(tf, [("Why this matters for the diagnostic   ", True, PURPLE_LT, 9.5),
          ("A maturity score that averages these ten gates hides the single "
           "one that will stop the programme. We report the index only "
           "alongside its bottleneck and its failed gates.",
           False, WHITE, 10.5)], after=0, line=1.2, first=True)

takeaway(sl, "This gate list is the spine of the proprietary diagnostic in "
             "section 04 — it is the part clients cannot self-assemble "
             "from open tooling.")
footer(sl)
notes(sl, "Emphasise ‘cannot be averaged away’. This is the single "
          "most common failure in agentic maturity models on the market, and "
          "it is our differentiator.")

# ========================================================== DIVIDER 02 ====
divider("02", "Value accrues to\ncontrol points",
        "Protocols are converging on open governance — and therefore "
        "commoditising.",
        ["Six architectural layers, and why no protocol spans them",
         "The nine-protocol register, verified as of 18 August 2026",
         "Where MCP fits — and the four jobs it is wrong for",
         "The nine scarce control points where value accrues",
         "Who gains and who loses as value migrates",
         "Twelve business-model archetypes, scored"])

# =========================================================== SLIDE 14 ====
sl = slide()
title_block(sl, "The architecture",
            "Agentic commerce needs six architectural layers — and no "
            "single protocol spans them")

layers = [
    ("6", "Operations and governance", "Observability, evaluation, audit, "
     "incident response, policy", "OpenTelemetry GenAI · OWASP · "
     "NIST AI RMF", PURPLE_DEEP),
    ("5", "Trust and transaction", "Delegated authority, mandates, payment, "
     "settlement, liability", "AP2 · Verifiable Intent · Trusted "
     "Agent · x402 · card rails", PURPLE_DARK),
    ("4", "Process and systems", "Systems of record, workflow, approval "
     "chains, fulfilment", "ERP · EDI · iPaaS · RPA",
     PURPLE),
    ("3", "Agent coordination", "Delegation, collaboration, multi-agent "
     "orchestration across boundaries", "A2A v1.0", PURPLE),
    ("2", "Capability exposure", "Tools, context, data access, capability "
     "invocation", "MCP 2026-07-28", PURPLE_LT),
    ("1", "Experience and access", "Assistants, storefronts, channels, agent "
     "discovery", "ACP · UCP · assistant surfaces", PURPLE_LT),
]
# vertical autonomy axis on the left
rect(sl, M, BODY_T, 0.34, 4.18, fill=PURPLE_LT, shape=MSO_SHAPE.UP_ARROW)
tf = vertical(box(sl, M + 0.42, BODY_T, 0.30, 4.18, anchor=MSO_ANCHOR.MIDDLE))
para(tf, "HIGHER AUTONOMY", size=8, bold=True, color=PURPLE_DARK,
     align=PP_ALIGN.CENTER, after=0, first=True)

BX = 1.40
BW = SW - M - BX
y = BODY_T
for num, name, desc, protos, col in layers:
    rect(sl, BX, y, BW, 0.62, fill=WHITE, line=GREY_LT, lw=0.75)
    rect(sl, BX, y, 0.58, 0.62, fill=col)
    tf = box(sl, BX, y + 0.17, 0.58, 0.3)
    para(tf, num, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
         after=0, first=True)
    tf = box(sl, BX + 0.72, y + 0.10, 2.55, 0.45, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, name, size=10.5, bold=True, color=INK, after=0, line=1.1,
         first=True)
    tf = box(sl, BX + 3.42, y + 0.10, 4.05, 0.45, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, desc, size=9, color=GREY, after=0, line=1.14, first=True)
    rect(sl, BX + 7.62, y + 0.13, 0.012, 0.36, fill=GREY_LT)
    tf = box(sl, BX + 7.82, y + 0.10, BW - 8.02, 0.45,
             anchor=MSO_ANCHOR.MIDDLE)
    para(tf, protos, size=9, bold=True, color=col if col != PURPLE_LT
         else PURPLE_DARK, after=0, line=1.14, first=True)
    y += 0.70

takeaway(sl, "Treating MCP as “the commerce stack” is a category "
             "error. It owns layer 2 decisively and touches nothing else.")
footer(sl)
notes(sl, "Autonomy increases as you go up the stack. Most client "
          "conversations start at layer 2 and have never considered layers 5 "
          "and 6 — which is exactly where our margin is.")

# =========================================================== SLIDE 15 ====
sl = slide()
title_block(sl, "The protocol register",
            "Nine mechanisms matter today — and every one of them is "
            "converging on open, foundation governance")

rows = [
    ["MCP", "2", "Agent-to-tool and data connectivity", "2026-07-28 stateless",
     "LF Projects / AAIF", "~½bn monthly SDK downloads; ~41% in production"],
    ["A2A", "3", "Cross-vendor agent communication", "v1.0, 9 Apr 2026",
     "Linux Foundation", "150+ orgs; Vertex, AgentCore, AI Foundry"],
    ["ACP", "5", "Agent-to-merchant checkout", "Beta, Apache 2.0",
     "OpenAI / Stripe / Meta", "PayPal, Agentforce — but scaled back Mar 2026"],
    ["UCP", "5", "Full shopping journey", "v2026-04-08",
     "Google + Shopify", "20+ endorsers; Walmart, Target, Etsy, Visa, Mastercard"],
    ["AP2", "5", "Delegated payment mandates", "v0.2, Apr 2026",
     "FIDO Alliance", "60+ partners; extends A2A and MCP"],
    ["Verifiable Intent", "5", "Tamper-proof log of authorised actions",
     "Contributed May 2026", "FIDO (Mastercard)", "Payments TWG chaired by "
     "Mastercard and Visa"],
    ["Trusted Agent Protocol", "5", "Distinguish legitimate agents from bots",
     "Open spec, Oct 2025", "Visa-led", "Cloudflare, Akamai, Adyen"],
    ["x402", "5", "HTTP-native payments", "Operational 14 Jul 2026",
     "x402 Fdn / Linux Fdn", "40 members — but volume heavily gamified"],
    ["REST, EDI, RPA, iPaaS", "2 / 4", "Existing integration", "Mature",
     "Various", "Often sufficient — and often the right answer"],
]
table(sl, M, BODY_T, CW,
      ["Protocol", "Layer", "Purpose", "Status", "Governance",
       "Adoption evidence"],
      rows, [1.72, 0.62, 2.75, 1.72, 1.72, 3.70], row_h=0.36, fsize=8.3,
      tint={0, 1})

takeaway(sl, "Standards position is worth having. Building a proprietary "
             "protocol is not — foundation governance will commoditise "
             "every one of these.")
footer(sl)
notes(sl, "Commerce semantics (ACP vs UCP) is the one contested layer. "
          "Everything else has consolidated. Advise clients to hedge across "
          "ACP and UCP rather than pick.")

# =========================================================== SLIDE 16 ====
sl = slide()
title_block(sl, "Protocol guidance",
            "MCP is the right tool for capability exposure — and the "
            "wrong tool for four other jobs")

quads = [("Where MCP is the best fit", PURPLE, PURPLE_MIST,
          ["Exposing enterprise tools and data to agents",
           "Internal agent-to-tool connectivity",
           "Anywhere an agent must discover and invoke a capability at runtime"]),
         ("Where MCP needs complements", PURPLE_DARK, WHITE,
          ["Cross-organisation coordination → A2A",
           "Commerce semantics → ACP or UCP",
           "Payment authority → AP2 / Verifiable Intent",
           "Settlement → x402 or card rails"]),
         ("Where existing APIs suffice", GREY, GREY_MIST,
          ["Deterministic, high-volume B2B transactions",
           "EDI and REST flows with no goal interpretation",
           "Anything that does not meet the agentic bar"]),
         ("Where MCP adds avoidable complexity", GREY, GREY_MIST,
          ["Simple request-response integrations",
           "Any integration with no agent in the loop",
           "Deterministic RPA dressed up as agentic"])]
x, y = M, BODY_T
for i, (head, col, bg, lines) in enumerate(quads):
    card(sl, x, y, 5.95, 1.72, "", head, lines, accent=col, fill=bg,
         border=GREY_LT, hsize=12, bsize=9.2)
    x = M + 6.28 if i % 2 == 0 else M
    if i % 2 == 1:
        y += 1.90

takeaway(sl, "Gartner projects 75% of API-gateway vendors will ship MCP "
             "features by end-2026 — the product opportunity is protocol "
             "gateways, adapters and a governed semantic layer above MCP.")
footer(sl)
notes(sl, "The two right-hand quadrants are the credibility move: we are the "
          "firm that tells clients when NOT to use the fashionable protocol.")

# =========================================================== SLIDE 17 ====
sl = slide()
title_block(sl, "The control points",
            "Connectivity will be commoditised — value accrues to nine "
            "scarce control points")

NW, NH = 2.62, 0.78
HX, HY, HW, HH = 5.29, 3.20, 2.75, 1.26          # hub
HCX, HCY = HX + HW / 2, HY + HH / 2

cps = [("01", "Customer and enterprise intent",       "Medium",    2.55, 1.72),
       ("02", "Proprietary domain data and knowledge", "High",     5.355, 1.72),
       ("03", "Differentiated executable capabilities", "High",    8.16, 1.72),
       ("04", "Identity, authority and policy",       "Very high", 0.55, 2.78),
       ("05", "Transaction and payment orchestration", "Very high", 10.16, 2.78),
       ("06", "Trusted distribution and ecosystem access", "Medium", 0.55, 4.42),
       ("07", "Workflow embedding and switching costs", "High",    10.16, 4.42),
       ("08", "Assurance, evaluation and accountability", "High",  2.55, 5.12),
       ("09", "Proprietary outcome data and learning loops", "High", 8.16, 5.12)]

# spokes first, so the nodes and the hub sit on top of them
for _, _, defens, nx, ny in cps:
    strong = defens == "Very high"
    spoke(sl, nx + NW / 2, ny + NH / 2, HCX, HCY,
          color=PURPLE if strong else GREY_LT, lw=1.5 if strong else 1.0)

for num, name, defens, nx, ny in cps:
    strong = defens == "Very high"
    rect(sl, nx, ny, NW, NH,
         fill=PURPLE_MIST if strong else WHITE,
         line=PURPLE if strong else GREY_LT, lw=1.25 if strong else 0.75,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.10)
    tf = box(sl, nx + 0.16, ny + 0.11, NW - 0.32, 0.18)
    rich(tf, [(num + "   ", True, PURPLE_LT, 7),
              (defens.upper(), True, PURPLE_DARK if strong else GREY, 7)],
         after=0, line=1.0, first=True)
    tf = box(sl, nx + 0.16, ny + 0.31, NW - 0.32, 0.42)
    para(tf, name, size=8.6, bold=True, color=INK, after=0, line=1.10,
         first=True)

rect(sl, HX, HY, HW, HH, fill=PURPLE_DEEP, shape=MSO_SHAPE.OVAL)
tf = box(sl, HX + 0.20, HY + 0.34, HW - 0.40, 0.60)
para(tf, "Durable", size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
     after=0, line=1.02, first=True)
para(tf, "value", size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
     after=0, line=1.02)

takeaway(sl, "Identity/authority and transaction orchestration are the two "
             "very-high-defensibility points — and both are exactly where "
             "the card networks and FIDO are consolidating.")
footer(sl)
notes(sl, "The ‘intent ownership’ control point (01) is the least "
          "proven: consumers are resisting delegation, so owning the surface "
          "has not converted into owning the transaction.")

# =========================================================== SLIDE 18 ====
sl = slide()
title_block(sl, "Value migration",
            "Value migrates from human-facing discovery to machine-readable "
            "capability, identity and settlement")

rect(sl, M, BODY_T, 5.95, 3.55, fill=PURPLE_MIST, line=PURPLE_LT, lw=0.9)
rect(sl, M, BODY_T, 5.95, 0.035, fill=PURPLE)
tf = box(sl, M + 0.24, BODY_T + 0.26, 5.5, 3.1)
para(tf, "LIKELY TO GAIN", size=10, bold=True, color=PURPLE_DARK, after=12,
     first=True)
for a, b in [("Payment networks",
              "They own settlement, credentialing, fraud — and now the "
              "trust layer"),
             ("Identity and trust providers",
              "FIDO, Okta, Cloudflare, Akamai: agent authentication and "
              "delegated-authority credentials"),
             ("Hyperscalers and cloud platforms",
              "They own the agent runtime and the gateway"),
             ("Well-integrated marketplaces",
              "Structured catalogues become the moat — Shopify Agentic "
              "Storefronts, Google Shopping Graph"),
             ("Data and knowledge providers",
              "Differentiated inputs become directly monetisable per call")]:
    rich(tf, [(a, True, INK, 10)], after=1, line=1.1)
    para(tf, b, size=8.8, color=GREY, after=9, line=1.16)

rect(sl, M + 6.28, BODY_T, 5.95, 3.55, fill=GREY_MIST, line=GREY_LT, lw=0.9)
rect(sl, M + 6.28, BODY_T, 5.95, 0.035, fill=GREY)
tf = box(sl, M + 6.52, BODY_T + 0.26, 5.5, 3.1)
para(tf, "LIKELY TO LOSE", size=10, bold=True, color=GREY, after=12,
     first=True)
for a, b in [("Pure discovery and SEO intermediaries",
              "The 8,000-vendor martech stack assumes buyers arrive at seller "
              "environments"),
             ("Online travel agencies",
              "Disintermediation risk is highest in travel distribution — "
              "OTA channel cost is 15–25% of room revenue"),
             ("Merchants with poor machine-readable data",
              "Agents route around unparseable suppliers — the penalty is "
              "silent and total"),
             ("Retail discovery advertising",
              "Retail media must be reinvented for a reader that is not human"),
             ("Anyone betting on a single protocol",
              "Foundation governance is commoditising the layer they hoped to "
              "own")]:
    rich(tf, [(a, True, INK, 10)], after=1, line=1.1)
    para(tf, b, size=8.8, color=GREY, after=9, line=1.16)

takeaway(sl, "The losing side is not losing to a competitor — it is "
             "losing to a reader that never visits. That is why data and "
             "capability readiness is the no-regret first move.")
footer(sl)
notes(sl, "Mirakl: businesses lose ~$15M a year to poor data quality. Under "
          "agentic procurement that compounds, because the agent simply never "
          "surfaces you.")

# =========================================================== SLIDE 19 ====
sl = slide()
title_block(sl, "Business models",
            "Twelve archetypes exist — but only five combine high "
            "defensibility with a Horizon-1 start")

rows = [
    ["Identity, authority and trust provider", "Enterprises and networks",
     "Per-credential, subscription", "Very high", "H1–H3"],
    ["Payment and settlement provider", "Merchants and agents",
     "Interchange, take-rate, FX", "Very high", "H1–H3"],
    ["Evaluation, security and assurance provider", "Enterprises",
     "Subscription, services", "High", "H1–H2"],
    ["Managed agent-operations provider", "Enterprises",
     "Managed fee, outcome-linked", "High", "H1–H2"],
    ["Orchestration platform", "Enterprises", "Consumption", "High",
     "H1–H3"],
    ["Enterprise agent provider", "Enterprises", "Licence, seat, usage",
     "Medium", "H1–H2"],
    ["Merchant capability provider", "Merchants", "Platform fee", "Medium",
     "H1–H2"],
    ["Data and knowledge provider", "Agents and enterprises",
     "Access licensing, x402 per call", "High", "H2–H3"],
    ["Demand-side personal agent", "Consumers", "Subscription, take-rate",
     "Low — assistant-dependent", "H2–H3"],
    ["Agent marketplace or broker", "Both sides", "Take-rate, listing",
     "High if liquidity achieved", "H3"],
    ["Outcome-based autonomous operator", "Enterprises", "Success fee",
     "High but capital-intensive", "H3"],
    ["Machine-resource marketplace", "Machines and agents", "Micro take-rate",
     "Medium", "H3"],
]
table(sl, M, BODY_T, CW,
      ["Archetype", "Payer / beneficiary", "Monetisation",
       "Defensibility", "Horizon"],
      rows, [3.90, 2.35, 2.90, 2.20, 0.98], row_h=0.30, fsize=8.4,
      tint={0, 1, 2, 3, 4})

takeaway(sl, "The five highlighted archetypes are where the firm should play. "
             "Four of the five are services businesses — which is the "
             "one place we start with an advantage.")
footer(sl)
notes(sl, "Classified on four non-overlapping axes: role in value chain, "
          "payer, monetisation, and control point. Ordered here by "
          "defensibility x horizon, not by size.")

# ========================================================== DIVIDER 03 ====
divider("03", "Where it lands,\nand when",
        "Three horizons, six sectors, one ranking.",
        ["The three-horizon interpretation — scenarios, not forecasts",
         "Six sectors ranked on Horizon-1 readiness",
         "Financial services: the most production-ready sector",
         "Retail: the most hyped and the least mature",
         "B2B and manufacturing: hardest, largest, data-gated",
         "Professional services: our own disruption"])

# =========================================================== SLIDE 21 ====
sl = slide()
title_block(sl, "The three horizons",
            "The restructuring is a three-to-seven year phenomenon — the "
            "next 18 months are foundations and bounded pilots")

rect(sl, M + 0.4, BODY_T + 0.30, CW - 0.8, 0.030, fill=INK)
for xpos, lab in [(M + 0.4, "Now"), (M + 4.30, "18 months"),
                  (M + 8.20, "36 months"), (M + 11.75, "7 years")]:
    rect(sl, xpos - 0.055, BODY_T + 0.20, 0.14, 0.14, fill=PURPLE,
         shape=MSO_SHAPE.OVAL)
    tf = box(sl, xpos - 0.75, BODY_T - 0.02, 1.5, 0.22)
    para(tf, lab, size=8.5, bold=True, color=GREY, align=PP_ALIGN.CENTER,
         after=0, first=True)

hz = [("Horizon 1", "0–18 months", "Readiness", PURPLE_DEEP,
       ["Bounded pilots with human-in-the-loop",
        "Internal workflow automation",
        "Machine-readable capability exposure",
        "Protocol experimentation",
        "Data and catalogue preparation",
        "Security, evaluation and observability",
        "Unit-level value evidence"]),
      ("Horizon 2", "18–36 months", "Scale", PURPLE,
       ["Scaled production deployment",
        "Cross-enterprise A2A coordination",
        "Agent-enabled channels (UCP / ACP)",
        "Delegated transactional authority",
        "Operating-model redesign",
        "Partner and platform restructuring",
        "Outcome and usage-based monetisation"]),
      ("Horizon 3", "3–7 years", "Frontier", PURPLE_LT,
       ["Persistent autonomous agents",
        "Machine-to-machine purchasing",
        "Agent marketplaces and brokers",
        "Dynamic contracting",
        "Machine-speed micropayments",
        "New economic actors",
        "Ecosystem restructuring"])]
x = M
for name, span, mode, col, items in hz:
    rect(sl, x, BODY_T + 0.62, 3.94, 3.02, fill=WHITE, line=GREY_LT, lw=0.75)
    rect(sl, x, BODY_T + 0.62, 3.94, 0.62, fill=col)
    tf = box(sl, x + 0.20, BODY_T + 0.72, 3.5, 0.24)
    para(tf, name.upper() + "   " + span, size=8.5, bold=True,
         color=WHITE if col != PURPLE_LT else PURPLE_DEEP, after=0, first=True)
    tf = box(sl, x + 0.20, BODY_T + 0.96, 3.5, 0.26)
    para(tf, mode, size=13, bold=True,
         color=WHITE if col != PURPLE_LT else PURPLE_DEEP, after=0, first=True)
    tf = box(sl, x + 0.20, BODY_T + 1.40, 3.55, 2.1)
    for i, it in enumerate(items):
        bullet(tf, it, size=9, color=C(0x3A, 0x3A, 0x46), mark=col,
               after=5, line=1.16, first=(i == 0))
    x += 4.14

takeaway(sl, "McKinsey’s $3–5 trillion figure is a Horizon-3 "
             "scenario, not a Horizon-1 fact. Forecasts across sources range "
             "35× — treat every one of them as a scenario to test.")
footer(sl)
notes(sl, "The 35x forecast range (eMarketer ~$144bn to McKinsey ~$5tn by "
          "2030) is driven purely by definitional boundaries. Never quote a "
          "market size without its definition.")

# =========================================================== SLIDE 22 ====
sl = slide()
title_block(sl, "Sector readiness",
            "Six sectors, one ranking: financial services is ready now; retail "
            "is the most hyped and the least mature")

rows = [
    ["1", "Financial services, payments, insurance", "High — production",
     "Credit, financial crime, claims, service", "Governance-first deployment "
     "and the settlement rails"],
    ["2", "Manufacturing and B2B distribution", "Medium-high",
     "Procurement, quote-to-cash, PO matching", "Machine-readable supplier and "
     "product data"],
    ["3", "Travel and hospitality", "Medium-high",
     "Booking, rebooking, disruption recovery", "Real-time machine-readable "
     "availability; direct-booking recapture"],
    ["4", "Professional services", "Medium",
     "Research, drafting, review, diligence", "Outcome-based pricing "
     "capability and delivery IP"],
    ["5", "Telecommunications and technology", "Medium — operations only",
     "Network operations, provisioning, NOC", "Exposing connectivity, identity "
     "and billing as agent-callable"],
    ["6", "FMCG, consumer products, retail", "Low on transaction",
     "Discovery and comparison only", "Catalogue normalisation, retail-media "
     "reinvention, identity linking"],
]
table(sl, M, BODY_T, CW,
      ["Rank", "Sector", "Horizon-1 readiness", "Where the value is now",
       "The control point to hold"],
      rows, [0.62, 3.15, 2.05, 3.10, 3.41], row_h=0.52, fsize=8.6,
      tint={0})

takeaway(sl, "Readiness tracks one variable: whether a human already sits on "
             "the authorisation step. Sectors that already have approval "
             "chains are ready; sectors that rely on impulse are not.")
footer(sl)
notes(sl, "This ranking is the sequencing logic for go-to-market. Lead with "
          "financial services and B2B; lead retail conversations with data "
          "readiness, never with checkout.")

# =========================================================== SLIDE 23 ====
sl = slide()
title_block(sl, "Sector spotlight: financial services",
            "Financial services is the most production-ready sector — and "
            "governance-first is the winning posture")

stats = [("70%", "of banking leaders are deploying or exploring agentic AI",
          "Sector survey, 2026"),
         ("30–50%", "workload reduction reported in banking operations",
          "Cost-to-serve is the primary driver"),
         ("20–60%", "productivity gain on credit memos, with 30% faster "
          "turnaround", "McKinsey, US bank"),
         ("80%", "faster insurance claims handling; near-instant settlements "
          "cited", "Allianz")]
x = M
for fig, lab, src in stats:
    statile(sl, x, BODY_T, 2.92, 1.78, fig, lab, src)
    x += 3.11

card(sl, M, BODY_T + 2.00, 5.95, 1.64, "Platform consolidation",
     "The core providers have already shipped",
     ["Fiserv agentOS; FIS with Anthropic for financial crime",
      "Oracle Financial Services agentic platform, extended to corporate "
      "banking April 2026",
      "Wolters Kluwer: 44% of finance teams expected to use agentic AI in 2026"],
     accent=PURPLE, hsize=11.5, bsize=8.8)

card(sl, M + 6.28, BODY_T + 2.00, 5.95, 1.64, "The rails land-grab",
     "Payments is the epicentre",
     ["Visa Intelligent Commerce Connect; Mastercard Agent Pay and AP4M",
      "Mastercard’s BVNK acquisition closed 3 August 2026 for up to $1.8bn",
      "Nearly half of banks and insurers are creating dedicated "
      "agent-supervision roles"],
     accent=PURPLE_DEEP, fill=PURPLE_MIST, hsize=11.5, bsize=8.8)

takeaway(sl, "“Compliance as a design directive” is the sector’s "
             "own language for our assurance proposition — sell it in "
             "their words, not ours.")
footer(sl)
notes(sl, "APAC is materially ahead of Western consumer deployment: Mastercard "
          "Agent Pay achieved live authenticated agentic transactions in Hong "
          "Kong (Mar 2026) and Thailand (Apr 2026).")

# =========================================================== SLIDE 24 ====
sl = slide()
title_block(sl, "Sector spotlight: retail and consumer",
            "Retail is the most hyped and the least mature — the durable "
            "pattern is “discover in AI, buy on site”")

card(sl, M, BODY_T, 3.94, 2.20, "The reality",
     "Discovery adoption is high; transaction adoption is not",
     ["A majority of consumers now use AI in shopping research",
      "Only 17% are comfortable completing the purchase",
      "OpenAI retreat; Walmart’s ~3× worse in-chat conversion"],
     accent=GREY, fill=GREY_MIST, hsize=11.5, bsize=9)

card(sl, M + 4.14, BODY_T, 3.94, 2.20, "The control points",
     "Three things worth owning",
     ["Machine-readable product data — catalogue normalisation is the moat",
      "Retail media reinvented for a non-human reader",
      "Loyalty and identity linking — UCP added Identity Linking in "
      "March 2026"],
     accent=PURPLE, fill=PURPLE_MIST, hsize=11.5, bsize=9)

card(sl, M + 8.28, BODY_T, 3.95, 2.20, "The likely winners",
     "Structured catalogues, not big brands",
     ["Shopify Agentic Storefronts activate UCP and native MCP by default",
      "Google Shopping Graph plus Merchant Center plus Google Pay plus UCP",
      "This combination is the live closed-platform risk"],
     accent=PURPLE_DEEP, hsize=11.5, bsize=9)

rect(sl, M, BODY_T + 2.42, CW, 1.22, fill=WHITE, line=GREY_LT, lw=0.75)
tf = box(sl, M + 0.26, BODY_T + 2.58, CW - 0.5, 0.9)
para(tf, "THE SEQUENCING WE RECOMMEND", size=9, bold=True, color=PURPLE_DARK,
     after=8, first=True)
rich(tf, [("Horizon 1  ", True, PURPLE, 9.5),
          ("Data and catalogue readiness; agent discoverability.    ",
           False, INK, 9.5),
          ("Horizon 2  ", True, PURPLE, 9.5),
          ("Agent-enabled channels; participate in UCP and ACP without "
           "betting on either.    ", False, INK, 9.5),
          ("Horizon 3  ", True, PURPLE, 9.5),
          ("Autonomous replenishment.", False, INK, 9.5)],
     after=0, line=1.35)

takeaway(sl, "Never open a retail conversation with checkout. Open it with "
             "catalogue readiness — it is the prerequisite for every "
             "later move and it pays back regardless of how autonomy lands.")
footer(sl)
notes(sl, "The closed-platform scenario is most live here. Google holds "
          "surface, catalogue, payment and protocol simultaneously.")

# =========================================================== SLIDE 25 ====
sl = slide()
title_block(sl, "Sector spotlight: B2B and manufacturing",
            "B2B is structurally the hardest and commercially the largest "
            "— and data quality is the competitive weapon")

card(sl, M, BODY_T, 5.95, 1.72, "The size of the prize",
     "The largest single value pool in the research",
     ["Gartner: by 2028, 90% of B2B buying will be AI-agent intermediated, "
      "pushing over $15 trillion of B2B spend through agent exchanges",
      "Forrester: ~20% of B2B sellers will face agent-led quote negotiations "
      "during 2026",
      "G2: 51% of B2B software buyers now begin in an AI chatbot; 69% switched "
      "vendor on chatbot guidance"], accent=PURPLE, hsize=11.5, bsize=8.8)

card(sl, M + 6.28, BODY_T, 5.95, 1.72, "Why it is hard",
     "Negotiated contracts, approval chains, EDI legacy",
     ["Autonomous transaction is structurally hardest where terms are "
      "negotiated rather than listed",
      "Existing EDI and REST flows are often sufficient — and MCP would "
      "add avoidable complexity",
      "Oracle a Leader in the 2026 Gartner MQ for Source-to-Pay; SAP Catalog "
      "Optimization Agent"], accent=PURPLE_DEEP, hsize=11.5, bsize=8.8)

rect(sl, M, BODY_T + 1.94, 5.95, 1.70, fill=BLACK)
rect(sl, M, BODY_T + 1.94, 0.045, 1.70, fill=PURPLE)
tf = box(sl, M + 0.30, BODY_T + 2.12, 5.4, 1.4)
para(tf, "THE COMPETITIVE WEAPON", size=8.5, bold=True, color=PURPLE, after=8,
     first=True)
rich(tf, [("$15M", True, WHITE, 26)], after=2, line=1.0)
para(tf, "average annual loss from poor data quality (Mirakl) — and it "
         "compounds under agentic procurement, because agents simply route "
         "around unparseable suppliers. The penalty is silent: you are not "
         "outbid, you are never surfaced.",
     size=9.5, color=C(0xB4, 0xB4, 0xBE), after=0, line=1.24)

card(sl, M + 6.28, BODY_T + 1.94, 5.95, 1.70, "The autonomy progression",
     "Earn each step — do not skip one",
     ["1. Research and market scanning",
      "2. Shortlisting and recommendation",
      "3. Approval-based ordering inside existing contracts",
      "4. Bounded negotiation within pre-agreed parameters"],
     accent=PURPLE, fill=PURPLE_MIST, hsize=11.5, bsize=9)

takeaway(sl, "Sell data readiness as revenue protection, not as an IT "
             "programme. In an agent-mediated market, unparseable is "
             "indistinguishable from absent.")
footer(sl)
notes(sl, "Gartner's 90% / $15tn is a 2028 prediction, not a measurement. Use "
          "it to size ambition, never to justify a business case.")

# =========================================================== SLIDE 26 ====
sl = slide()
title_block(sl, "Sector spotlight: professional services",
            "Our own industry is the clearest business-model disruption — "
            "and outcome-based pricing is arriving more slowly than claimed")

rect(sl, M, BODY_T, 5.95, 3.55, fill=PURPLE_MIST, line=PURPLE_LT, lw=0.9)
rect(sl, M, BODY_T, 5.95, 0.035, fill=PURPLE)
tf = box(sl, M + 0.24, BODY_T + 0.26, 5.5, 3.1)
para(tf, "THE SIGNAL", size=9.5, bold=True, color=PURPLE_DARK, after=11,
     first=True)
for a, b in [("McKinsey", "~25% of global fees now outcome-linked; "
              "“Lilli” deployed at scale"),
             ("BCG", "expects AI-tied work at ~40% of revenue in 2026, "
              "up from ~20% in 2024"),
             ("Bain", "~30% today, heading toward 50%"),
             ("Accenture", "$3.6bn AI bookings in FY2025 — up 120% "
              "year on year"),
             ("EY and Deloitte", "have both published the accounting "
              "playbooks: EY on outcome-based pricing under ASC 606 "
              "(17 Feb 2026); Deloitte DART on agentic-AI pricing "
              "(4 Jun 2026)"),
             ("Legal", "Harvey at 1,300+ organisations and 100,000+ lawyers; "
              "LexisNexis Protégé GA with 300+ workflows")]:
    rich(tf, [(a + "  ", True, INK, 9.5), (b, False, GREY, 9.5)],
         after=8, line=1.20)

rect(sl, M + 6.28, BODY_T, 5.95, 3.55, fill=GREY_MIST, line=GREY_LT, lw=0.9)
rect(sl, M + 6.28, BODY_T, 5.95, 0.035, fill=GREY)
tf = box(sl, M + 6.52, BODY_T + 0.26, 5.5, 3.1)
para(tf, "THE COUNTER-SIGNAL", size=9.5, bold=True, color=GREY, after=11,
     first=True)
bullet(tf, "The shift is sluggish. Even at McKinsey, roughly 75% of fees "
           "remain traditional billing.", size=9.8, mark=GREY, after=9,
       line=1.20)
bullet(tf, "No public data quantifies what share of contracts are actually "
           "being renegotiated to outcomes — the headline percentages "
           "describe intent, not executed contracts.", size=9.8, mark=GREY,
       after=9, line=1.20)
bullet(tf, "Outcome pricing transfers delivery risk to us. Without the "
           "assurance and evaluation capability from section 01, that is an "
           "uninsured transfer.", size=9.8, mark=GREY, after=9, line=1.20)
bullet(tf, "The accounting treatment is only twelve months old — revenue "
           "recognition is a live constraint, not a solved one.",
       size=9.8, mark=GREY, after=0, line=1.20)

takeaway(sl, "Treat outcome-based pricing as a capability to build "
             "deliberately over Horizon 2 — not as a forced pivot. Build "
             "the assurance capability first; it is what makes the risk "
             "transfer survivable.")
footer(sl)
notes(sl, "This is the slide that makes the deck about us, not about clients. "
          "The counter-signal column is what stops an over-reaction in the "
          "room.")

# ========================================================== DIVIDER 04 ====
divider("04", "What we do\nabout it",
        "Three theses, kept deliberately separate.",
        ["Why conflating the three theses is the common strategic error",
         "What we advise clients to fund — and to defer",
         "What the firm should build, and refuse to build",
         "Where the firm should invest externally",
         "The three-horizon roadmap, with two decision gates"])

# =========================================================== SLIDE 28 ====
sl = slide()
title_block(sl, "Three theses",
            "Three investment theses answer three different questions — "
            "conflating them is the most common strategic error")

th = [("A", "The client thesis", "What should our clients fund?",
       ["Owner: industry and client-account leadership",
        "Currency: client capital and client risk appetite",
        "Failure mode: recommending what we happen to sell"],
       PURPLE_DEEP),
      ("B", "The firm-capability thesis", "What should we build ourselves?",
       ["Owner: practice and capability leadership",
        "Currency: our own investment budget and partner time",
        "Failure mode: building what open tooling will commoditise"],
       PURPLE),
      ("C", "The market thesis", "Where should we invest, partner or acquire?",
       ["Owner: corporate development and alliances",
        "Currency: balance sheet and alliance slots",
        "Failure mode: buying into a category we could simply partner in"],
       PURPLE_LT)]
x = M
for letter, name, q, lines, col in th:
    rect(sl, x, BODY_T, 3.94, 3.30, fill=WHITE, line=GREY_LT, lw=0.75)
    rect(sl, x, BODY_T, 3.94, 0.70, fill=col)
    tf = box(sl, x + 0.22, BODY_T + 0.14, 3.5, 0.45)
    rich(tf, [(letter + "   ", True, WHITE if col != PURPLE_LT
               else PURPLE_DEEP, 18),
              (name, True, WHITE if col != PURPLE_LT else PURPLE_DEEP, 12.5)],
         after=0, line=1.0, first=True)
    tf = box(sl, x + 0.22, BODY_T + 0.92, 3.5, 0.5)
    para(tf, q, size=11.5, bold=True, color=INK, after=0, line=1.12,
         first=True)
    tf = box(sl, x + 0.22, BODY_T + 1.62, 3.5, 1.5)
    for i, ln in enumerate(lines):
        bullet(tf, ln, size=9.2, color=GREY, mark=col if col != PURPLE_LT
               else PURPLE, after=7, line=1.18, first=(i == 0))
    x += 4.14

takeaway(sl, "Rank these three separately and review them separately. A "
             "category can be a strong client recommendation and a poor firm "
             "investment at the same time — observability is exactly that.")
footer(sl)
notes(sl, "Observability is the worked example: we should tell every client to "
          "invest in it, we should never build it, and we should partner for "
          "it. One category, three different answers.")

# =========================================================== SLIDE 29 ====
sl = slide()
title_block(sl, "Thesis A — the client thesis",
            "Advise clients to fund three no-regret foundations now, and to "
            "defer consumer autonomous checkout")

rows = [
    ["Machine-readable data and capability exposure",
     "No-regret foundation", "H1",
     "Prerequisite for every downstream use case — agents route around "
     "unparseable data"],
    ["Identity, delegated authority and consent architecture",
     "No-regret foundation", "H1",
     "The binding gate for any autonomous execution; the regulatory direction "
     "of travel (FIDO, HM Treasury)"],
    ["Agent evaluation, observability and security",
     "No-regret foundation", "H1",
     "88% incident rate makes this a production blocker, not a maturity nicety"],
    ["Bounded enterprise use cases with human-in-the-loop",
     "Near-term priority", "H1",
     "Quantified ROI at low exposure — procurement, PO matching, "
     "financial crime, service"],
    ["UCP / ACP participation and card-network agentic rails",
     "Near-term priority", "H2",
     "Preserves channel access; hedge across both rather than picking a winner"],
    ["Outcome-based and autonomous-operator models",
     "Strategic option", "H2–H3",
     "High value but unproven economics and unresolved liability"],
    ["Machine-native settlement (x402)", "Monitor", "H3",
     "Volume is largely gamified today — watch for genuine adoption"],
    ["Building proprietary protocols", "Avoid", "—",
     "Protocols are commoditising under foundation governance"],
]
table(sl, M, BODY_T, CW,
      ["Recommendation", "Classification", "Horizon", "Rationale"],
      rows, [4.15, 2.05, 0.85, 5.18], row_h=0.38, fsize=8.6,
      chip_col=1, chip_map=CHIPS)

rect(sl, M, BODY_T + 3.62, CW, 0.46, fill=BLACK)
tf = box(sl, M + 0.26, BODY_T + 3.70, CW - 0.5, 0.3, anchor=MSO_ANCHOR.MIDDLE)
rich(tf, [("Default posture   ", True, PURPLE, 9),
          ("BUILD", True, WHITE, 9.5),
          (" data, capability and identity foundations ·  ", False,
           C(0xB4, 0xB4, 0xBE), 9.5),
          ("BUY or PARTNER", True, WHITE, 9.5),
          (" for observability, security and payment rails ·  ", False,
           C(0xB4, 0xB4, 0xBE), 9.5),
          ("DEFER", True, WHITE, 9.5),
          (" consumer autonomous checkout", False, C(0xB4, 0xB4, 0xBE), 9.5)],
     after=0, line=1.1, first=True)

takeaway(sl, "The three no-regret foundations pay back under every scenario "
             "in section 05 — including the ones where this thesis is "
             "wrong. That is what makes them the opening recommendation.")
footer(sl)
notes(sl, "If challenged on any single recommendation, retreat to the "
          "no-regret test: does it pay back even if autonomy never arrives? "
          "For the top three, it does.")

# =========================================================== SLIDE 30 ====
sl = slide()
title_block(sl, "Thesis B — the firm-capability thesis",
            "The firm should build proprietary diagnostics and assurance IP "
            "— and should refuse to build an observability platform")

rows = [
    ["Proprietary agentic-commerce maturity diagnostic",
     "No-regret foundation",
     "Eight dimensions, ten gates, sector benchmarks; a 64–96 question "
     "bank",
     "Recurring assessment revenue; lead-gen into transformation"],
    ["Trust and control library plus assurance methodology",
     "No-regret foundation",
     "Six control domains mapped to a risk-control matrix",
     "High willingness-to-pay given the 88% incident rate"],
    ["Sector benchmarks and use-case libraries", "Near-term priority",
     "Six priority sectors, reusable across engagements",
     "Differentiation and delivery speed"],
    ["Protocol accelerators, reusable MCP assets and adapters",
     "Near-term priority",
     "Gateways, adapters and a governed semantic layer over MCP",
     "Delivery speed plus recurring maintenance revenue"],
    ["Evaluation and observability accelerators", "Near-term priority",
     "Assembled on OpenTelemetry, Confident AI and Arthur — not built",
     "Partner-led and fast to stand up"],
    ["Managed agent-operations and assurance service", "Strategic option",
     "Monitoring, evaluation and incident response as a service",
     "Recurring-revenue annuity; sticky by design"],
    ["Acquisition of niche agent-security or evaluation boutiques",
     "Strategic option", "A fragmented market with real talent scarcity",
     "Talent and IP acceleration"],
    ["A full proprietary observability platform", "Avoid",
     "Open tooling is moving too fast to out-run",
     "Partner instead — the category will commoditise"],
]
table(sl, M, BODY_T, CW,
      ["Investment", "Classification", "What it is", "Commercial model"],
      rows, [3.85, 2.05, 3.35, 2.98], row_h=0.42, fsize=8.5,
      chip_col=1, chip_map=CHIPS)

takeaway(sl, "The strategy is assemble-then-build: wrap best-of-breed open "
             "tooling inside proprietary diagnostic and control IP, and build "
             "only the executive-assessment, prioritisation and benchmark "
             "layers where the IP is genuinely defensible.")
footer(sl)
notes(sl, "The market gap is real: no single tool today combines executive "
          "assessment, automated architecture discovery, data readiness, agent "
          "evaluation, security gates, prioritisation and sector benchmarks. "
          "That gap is the product.")

# =========================================================== SLIDE 31 ====
sl = slide()
title_block(sl, "Thesis C — the market thesis",
            "Externally, three categories are near-term priorities — and "
            "one is an explicit avoid")

cats = [("Identity, delegated authority and trust for agents",
         "Near-term priority",
         "Very high defensibility, standards-anchored at FIDO, high "
         "willingness-to-pay"),
        ("Agent security, assurance and evaluation", "Near-term priority",
         "High willingness-to-pay in a fragmented market with credible "
         "acquisition targets"),
        ("Data readiness and machine-readable catalogue services",
         "Near-term priority",
         "Clear and quantified pain; buyers already understand the problem"),
        ("Protocol gateways, adapters and governed semantic layers",
         "Strategic option",
         "Durable precisely if standards fragment — a hedge against our "
         "own convergence call"),
        ("Managed agent operations", "Strategic option",
         "Annuity revenue, but dependent on proving assurance demand first"),
        ("Machine-native settlement and agent marketplaces", "Frontier bet",
         "Large addressable market but early, gamified volume and "
         "platform-dependent"),
        ("Consumer personal-agent applications", "Avoid / monitor",
         "Low defensibility and wholly assistant-dependent")]
y = BODY_T
for name, cls, why in cats:
    fill, txt = CHIPS[cls]
    rect(sl, M, y, CW, 0.52, fill=WHITE, line=GREY_LT, lw=0.75)
    rect(sl, M, y, 0.045, 0.52, fill=fill if cls != "Avoid / monitor"
         else GREY)
    tf = box(sl, M + 0.24, y + 0.07, 4.55, 0.38, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, name, size=10.5, bold=True, color=INK, after=0, line=1.1,
         first=True)
    rect(sl, M + 5.00, y + 0.11, 1.72, 0.30, fill=fill)
    tf = box(sl, M + 5.00, y + 0.17, 1.72, 0.2)
    para(tf, cls.upper(), size=7, bold=True, color=txt,
         align=PP_ALIGN.CENTER, after=0, first=True)
    tf = box(sl, M + 6.95, y + 0.07, CW - 7.15, 0.38, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, why, size=9, color=GREY, after=0, line=1.16, first=True)
    y += 0.585

takeaway(sl, "Scored on addressable value pool, willingness-to-pay, horizon, "
             "defensibility, switching costs, capital intensity, "
             "technical and regulatory risk, competitive intensity and "
             "platform dependence.")
footer(sl)
notes(sl, "Keep this ranking separate from thesis A and B. Consumer "
          "personal-agent apps are an avoid for our balance sheet even though "
          "they are a legitimate thing for a client to experiment with.")

# =========================================================== SLIDE 32 ====
sl = slide()
title_block(sl, "The roadmap",
            "Two decision gates govern the roadmap — at 12 months on "
            "assurance demand, and at 36 months on outcome economics")

hz = [("Horizon 1", "0–18 months", PURPLE_DEEP,
       ["Build the maturity diagnostic, trust/control library and sector "
        "use-case libraries",
        "Stand up an agent-security and assurance practice — partner on "
        "tooling",
        "Launch data/capability readiness and MCP-inventory assessments",
        "Establish alliances with a hyperscaler, a card network and an "
        "identity provider",
        "Recruit and retain agent-security and integration talent"],
       "Diagnostic engagements sold · assurance ARR · alliance "
       "certifications",
       "GATE AT 12 MONTHS — if assurance and diagnostic demand is "
       "proven, expand into managed operations"),
      ("Horizon 2", "18–36 months", PURPLE,
       ["Scale managed agent-operations as recurring revenue",
        "Deliver cross-enterprise A2A workflow transformations",
        "Build outcome-based pricing capability internally, learning from the "
        "EY, Deloitte and BCG playbooks",
        "Develop sector reference architectures"],
       "Managed-service ARR · outcome-based revenue share · "
       "repeat-client rate",
       "GATE AT 36 MONTHS — if outcome-based economics prove out, shift "
       "more of the book to success-fee models"),
      ("Horizon 3", "3–7 years", PURPLE_LT,
       ["Position for agent-marketplace and machine-commerce advisory",
        "Develop dynamic-contracting IP",
        "Consider selective build or acquire in identity and assurance if the "
        "market consolidates"],
       "Share of revenue from machine-commerce categories · IP licensing",
       "")]
x = M
for name, span, col, acts, kpis, gate in hz:
    rect(sl, x, BODY_T, 3.94, 4.05, fill=WHITE, line=GREY_LT, lw=0.75)
    rect(sl, x, BODY_T, 3.94, 0.52, fill=col)
    tf = box(sl, x + 0.20, BODY_T + 0.13, 3.5, 0.3)
    rich(tf, [(name + "   ", True, WHITE if col != PURPLE_LT else PURPLE_DEEP,
               12),
              (span, False, WHITE if col != PURPLE_LT else PURPLE_DEEP, 9)],
         after=0, line=1.0, first=True)
    tf = box(sl, x + 0.20, BODY_T + 0.68, 3.55, 2.2)
    for i, a in enumerate(acts):
        bullet(tf, a, size=8.8, color=C(0x3A, 0x3A, 0x46), mark=col,
               after=6, line=1.18, first=(i == 0))
    rect(sl, x + 0.20, BODY_T + 2.92, 3.55, 0.012, fill=GREY_LT)
    tf = box(sl, x + 0.20, BODY_T + 3.02, 3.55, 0.55)
    para(tf, "KPIs", size=7.5, bold=True, color=PURPLE_DARK, after=3,
         first=True)
    para(tf, kpis, size=8.2, color=GREY, after=0, line=1.16)
    if gate:
        rect(sl, x + 0.20, BODY_T + 3.58, 3.55, 0.38, fill=BLACK)
        tf = box(sl, x + 0.30, BODY_T + 3.63, 3.35, 0.3,
                 anchor=MSO_ANCHOR.MIDDLE)
        para(tf, gate, size=7.2, bold=True, color=WHITE, after=0, line=1.12,
             first=True)
    x += 4.14

takeaway(sl, "Nothing in Horizon 2 is committed until the 12-month gate "
             "clears. The roadmap is deliberately optioned, because the "
             "falsifiers in section 05 are live.")
footer(sl)
notes(sl, "Present the gates as the point of the slide, not the activities. "
          "Executives should leave knowing what evidence unlocks the next "
          "tranche of investment.")

# ========================================================== DIVIDER 05 ====
divider("05", "What would prove\nus wrong",
        "The falsifiers and indicators, declared before the fact.",
        ["Five developments that would falsify or weaken the thesis",
         "Seven leading indicators, and the two that matter most",
         "The assumptions this research is built on",
         "Six decisions for the executive team"])

# =========================================================== SLIDE 34 ====
sl = slide()
title_block(sl, "Falsifiers",
            "Five developments would falsify or materially weaken this thesis "
            "— we are declaring them before the fact")

fals = [("01", "Consumer autonomous checkout recovers",
         "If conversion reaches parity with redirect, the consumer sub-thesis "
         "we downgrade here would need upgrading — and Horizon-1 spend "
         "should shift toward consumer channels."),
        ("02", "Protocols fragment permanently",
         "If ACP and UCP never converge, integration costs rise and the "
         "fragmentation counter-thesis wins. Gateways and adapters become "
         "more valuable, not less."),
        ("03", "Rails commoditise to near-zero margin",
         "If incumbents capture all the value, there is no new value pool and "
         "no consulting annuity — only a one-off integration wave."),
        ("04", "Regulatory clampdown stalls delegated authority",
         "EU high-risk classification of agentic commerce, or PSD3 SCA rules "
         "that block autonomous payment, would freeze Horizon 2."),
        ("05", "Machine-native settlement never becomes genuine",
         "Roughly half of x402 activity was wash or gamified at the March "
         "2026 snapshot, with only ~$28,000 a day genuine. If that never "
         "converts, Horizon 3 shrinks.")]
x, y = M, BODY_T
for i, (num, head, body) in enumerate(fals):
    w = 3.94 if i < 3 else 5.95
    if i == 3:
        x, y = M, BODY_T + 1.82
    rect(sl, x, y, w, 1.60, fill=WHITE, line=GREY_LT, lw=0.75)
    rect(sl, x, y, w, 0.035, fill=PURPLE)
    tf = box(sl, x + 0.20, y + 0.20, w - 0.4, 0.22)
    para(tf, num, size=9, bold=True, color=PURPLE_LT, after=0, first=True)
    tf = box(sl, x + 0.20, y + 0.44, w - 0.4, 0.44)
    para(tf, head, size=11, bold=True, color=INK, after=0, line=1.10,
         first=True)
    tf = box(sl, x + 0.20, y + 0.90, w - 0.4, 0.62)
    para(tf, body, size=8.5, color=GREY, after=0, line=1.18, first=True)
    x += (4.14 if i < 3 else 6.28)

takeaway(sl, "None of these five is currently triggered. Two of them — "
             "01 and 05 — are the ones we should instrument now, because "
             "they would change the investment case rather than merely delay "
             "it.")
footer(sl)
notes(sl, "Declaring falsifiers up front is the credibility move with an "
          "executive audience. It converts a point of view into a testable "
          "position.")

# =========================================================== SLIDE 35 ====
sl = slide()
title_block(sl, "Leading indicators",
            "Seven indicators tell us whether we are right — and two of "
            "them matter more than the other five")

inds = [("MCP production-adoption rate", "~41% today, with OAuth 2.1 at ~8.5%",
         False),
        ("ACP / UCP convergence signals",
         "The one genuinely contested layer in the stack", False),
        ("x402 genuine, wash-filtered volume",
         "The falsifier for the whole Horizon-3 case", True),
        ("Card-network agentic transaction counts in live APAC markets",
         "Hong Kong and Thailand are the leading edge", False),
        ("Consumer autonomous-checkout conversion versus redirect",
         "The single most consequential number in this deck", True),
        ("FIDO agentic-standards ratification",
         "Determines how fast delegated authority becomes routine", False),
        ("PSD3 / PSR application date and any agent-specific provisions",
         "Earliest ~H2 2027, potentially slipping to 2028", False)]
y = BODY_T
for name, why, key in inds:
    rect(sl, M, y, CW, 0.52, fill=PURPLE_MIST if key else WHITE,
         line=PURPLE_LT if key else GREY_LT, lw=0.9 if key else 0.75)
    rect(sl, M, y, 0.045, 0.52, fill=PURPLE if key else GREY_LT)
    tf = box(sl, M + 0.26, y + 0.07, 5.6, 0.38, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, name, size=10.5, bold=key, color=INK, after=0, line=1.1,
         first=True)
    tf = box(sl, M + 6.10, y + 0.07, CW - 7.6, 0.38, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, why, size=9, color=GREY, after=0, line=1.14, first=True)
    if key:
        rect(sl, SW - M - 1.32, y + 0.12, 1.20, 0.28, fill=PURPLE)
        tf = box(sl, SW - M - 1.32, y + 0.18, 1.20, 0.2)
        para(tf, "PRIORITY", size=7, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, after=0, first=True)
    y += 0.585

takeaway(sl, "Instrument the two priority indicators as a standing quarterly "
             "review item. They are the two datapoints that would most change "
             "the thesis — everything else is confirmation.")
footer(sl)
notes(sl, "Propose a named owner and a quarterly cadence in the room. An "
          "indicator with no owner is a footnote, not a control.")

# =========================================================== SLIDE 36 ====
sl = slide()
title_block(sl, "Evidence discipline",
            "Four assumptions shape every number in this deck — they "
            "should be challenged before the conclusions are")

asm = [("01", "The agentic bar is defined and enforced",
        "“Agentic” requires goal interpretation, decision-making, "
        "tool use, adaptation or delegated action. Deterministic RPA and EDI "
        "are excluded unless they meet that bar."),
       ("02", "Announcements are not adoption",
        "Pilots are not production. Every adoption figure in this deck "
        "distinguishes the two, and vendor claims are triangulated with "
        "independent evidence or flagged where they cannot be."),
       ("03", "Market sizes are scenarios, not predictions",
        "Forecasts range 35× across sources — from eMarketer’s "
        "~$144bn to McKinsey’s ~$5tn by 2030 — driven purely by "
        "definitional boundaries."),
       ("04", "Two variables carry most of the sensitivity",
        "The findings are most sensitive to whether consumer autonomous "
        "checkout conversion recovers, and whether machine-native settlement "
        "volume becomes genuine. Both are flagged as falsifiers.")]
x, y = M, BODY_T
for i, (num, head, body) in enumerate(asm):
    rect(sl, x, y, 5.95, 1.72, fill=WHITE, line=GREY_LT, lw=0.75)
    rect(sl, x, y, 5.95, 0.035, fill=PURPLE)
    tf = box(sl, x + 0.22, y + 0.22, 5.5, 0.22)
    para(tf, "ASSUMPTION " + num, size=8, bold=True, color=PURPLE_LT, after=0,
         first=True)
    tf = box(sl, x + 0.22, y + 0.48, 5.5, 0.32)
    para(tf, head, size=11.5, bold=True, color=INK, after=0, line=1.1,
         first=True)
    tf = box(sl, x + 0.22, y + 0.86, 5.5, 0.76)
    para(tf, body, size=9, color=GREY, after=0, line=1.20, first=True)
    x = M + 6.28 if i % 2 == 0 else M
    if i % 2 == 1:
        y += 1.90

takeaway(sl, "If the room disagrees with a conclusion, take the disagreement "
             "back to one of these four assumptions. That is where the real "
             "argument lives.")
footer(sl)
notes(sl, "Use this slide defensively when a senior stakeholder challenges a "
          "number. Almost every challenge resolves to assumption 01 (what "
          "counts as agentic) or 03 (market size definitions).")

# =========================================================== SLIDE 37 ====
sl = slide()
rect(sl, 0, 0, SW, SH, fill=BLACK)
rect(sl, 0, 0, SW, 0.10, fill=PURPLE)

tf = box(sl, M + 0.35, 0.62, 8.0, 0.26)
para(tf, "Executive conclusion", size=9, bold=True, color=PURPLE, after=0,
     first=True, caps=True)
tf = box(sl, M + 0.35, 0.98, 11.4, 0.62)
para(tf, "Six decisions for the executive team", size=25, bold=True,
     color=WHITE, after=0, line=1.05, first=True)
rect(sl, M + 0.35, 1.78, 1.15, 0.035, fill=PURPLE)

dec = [("Believe",
        "The front is real and consequential — but the next 18 months of "
        "value are in enterprise and B2B workflows and the "
        "identity-payment-trust rails, not autonomous consumer shopping."),
       ("Advise",
        "Fund machine-readable capability, identity and delegated authority, "
        "and agent assurance as no-regret foundations. Run bounded pilots with "
        "a human in the loop. Defer consumer autonomous checkout."),
       ("Build",
        "A proprietary maturity diagnostic, a trust/control and assurance "
        "library, sector use-case libraries and benchmarks, and protocol "
        "accelerators — assembling, not rebuilding, open tooling."),
       ("Partner for",
        "Observability and evaluation, payment rails and identity — via "
        "alliances with one hyperscaler, one card network and one identity "
        "provider."),
       ("Invest in",
        "Agent identity and trust, security and assurance, and data readiness "
        "near-term; managed operations and protocol gateways as options; "
        "machine-native settlement as a frontier bet."),
       ("Monitor",
        "The seven leading indicators — above all whether consumer "
        "autonomous-checkout conversion recovers, and whether machine-native "
        "settlement volume becomes genuine.")]
x, y = M + 0.35, 2.10
for i, (verb, body) in enumerate(dec):
    rect(sl, x, y, 5.72, 1.28, fill=C(0x12, 0x12, 0x18),
         line=C(0x2E, 0x2E, 0x38), lw=0.75)
    rect(sl, x, y, 0.04, 1.28, fill=PURPLE)
    tf = box(sl, x + 0.26, y + 0.18, 5.2, 0.26)
    para(tf, verb.upper(), size=9.5, bold=True, color=PURPLE, after=0,
         first=True)
    tf = box(sl, x + 0.26, y + 0.48, 5.2, 0.72)
    para(tf, body, size=9.2, color=C(0xC8, 0xC8, 0xD2), after=0, line=1.22,
         first=True)
    x = M + 6.62 if i % 2 == 0 else M + 0.35
    if i % 2 == 1:
        y += 1.40

tf = box(sl, M + 0.35, 6.72, 10.0, 0.24)
para(tf, "Prepared 18 August 2026 · Adoption figures distinguish "
         "announcements from production · Forecasts are labelled as "
         "scenarios · Vendor claims are triangulated or flagged",
     size=7.5, color=GREY, after=0, first=True)
_n[0] += 1
tf = box(sl, SW - M - 1.2, 6.72, 1.2, 0.24)
para(tf, str(_n[0]), size=7.5, bold=True, color=PURPLE,
     align=PP_ALIGN.RIGHT, after=0, first=True)
notes(sl, "Close on the two monitored variables. They are what turns this "
          "from a point of view into a managed position.")

# ---------------------------------------------------------------- save ----
here = os.path.dirname(os.path.abspath(__file__))
out = os.path.join(
    here, "2026.08.18-Machine-Mediated-Commerce-Front-Executive-Deck.pptx")
prs.save(out)
print("Wrote %s (%d slides)" % (out, len(prs.slides.__iter__.__self__._sldIdLst)))
